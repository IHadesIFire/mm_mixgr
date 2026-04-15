"""Generate a simple 2-slide progress PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]


def add_title(slide, text, top=0.3):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True


def add_text(slide, text, left, top, width, height, size=18, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
    return box


# ========== Slide 1: Reproduction ==========
s1 = prs.slides.add_slide(blank)
add_title(s1, "Slide 1 — Reproduction of MRMR Baseline (GME-Qwen2-VL-7B)")

add_text(s1, "Setup", 0.5, 1.2, 6, 0.5, size=22, bold=True)
add_text(
    s1,
    "Dataset: MRMR Knowledge subtask (ICLR 2026)\n"
    "  • 555 queries, 26,223 corpus docs\n"
    "  • 4 coarse domains: Art / Humanities / Medicine / Science\n"
    "Encoder: GME-Qwen2-VL-7B-Instruct (3584-dim, L2-normalized)\n"
    "Metric: nDCG@10",
    0.5, 1.7, 8, 2.5, size=16,
)

add_text(s1, "Result — matches paper", 7.0, 1.2, 6, 0.5, size=22, bold=True)

# simple 2-col table
rows = [
    ("Split",        "Ours",   "MRMR paper"),
    ("Overall",      "0.519",  "~0.519"),
    ("Art",          "0.563",  "—"),
    ("Humanities",   "0.469",  "—"),
    ("Medicine",     "0.484",  "—"),
    ("Science",      "0.545",  "—"),
]
from pptx.util import Inches
table = s1.shapes.add_table(len(rows), 3, Inches(7.0), Inches(1.8), Inches(5.8), Inches(3.2)).table
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            if i == 0:
                p.font.bold = True

add_text(
    s1,
    "Conclusion: overall nDCG@10 reproduces the paper's reported number → baseline pipeline is correct.",
    0.5, 6.3, 12.3, 0.8, size=18, bold=True,
)


# ========== Slide 2: Badcase & next step ==========
s2 = prs.slides.add_slide(blank)
add_title(s2, "Slide 2 — Bad Case Analysis: Three Failure Modes")

add_text(s2, "162 / 555 queries have nDCG@10 = 0 (29%). Diagnosed via sentence-level similarity:",
         0.5, 1.1, 12, 0.6, size=17)

add_text(s2, "A. Dilution", 0.5, 1.9, 4, 0.5, size=20, bold=True)
add_text(
    s2,
    "Gold doc contains a sentence that matches the query well,\n"
    "but whole-doc pooling drowns the signal.\n"
    "→ proposition-level retrieval should help.",
    0.5, 2.5, 4.1, 2.0, size=14,
)

add_text(s2, "B. Misalignment", 4.7, 1.9, 4, 0.5, size=20, bold=True)
add_text(
    s2,
    "No sentence in the gold doc aligns with the query\n"
    "semantically / lexically.\n"
    "→ query rewrite / HyDE territory.",
    4.7, 2.5, 4.1, 2.0, size=14,
)

add_text(s2, "C. Noise", 8.9, 1.9, 4, 0.5, size=20, bold=True)
add_text(
    s2,
    "Best gold sentence has very low absolute similarity\n"
    "to the query (< 0.35).\n"
    "→ qrel label likely noisy.",
    8.9, 2.5, 4.1, 2.0, size=14,
)

add_text(s2, "Preliminary distribution (20 zero-nDCG queries)", 0.5, 4.7, 8, 0.5, size=18, bold=True)
add_text(
    s2,
    "• B_misalign: 65%    • A_dilution: 30%    • other: 5%    • C_noise: 0%\n"
    "Per-domain: Art → all A;  Medicine → mostly B;  Science → mixed.",
    0.5, 5.3, 12, 1.2, size=15,
)

add_text(s2, "What we are doing now", 0.5, 6.3, 6, 0.5, size=18, bold=True)
add_text(
    s2,
    "1. Running full 162-query dilution diagnosis to confirm A/B/C split.\n"
    "2. Designing per-domain remedy: proposition-level retrieval for A, query rewrite for B.",
    0.5, 6.8, 12.3, 0.8, size=14,
)

out = "progress.pptx"
prs.save(out)
print(f"[ok] saved {out}")
